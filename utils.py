import calendar
import datetime
import math
from io import BytesIO

import df2img
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt

# Daily Report - A link to download the daily report in pptx format but this is still on dev.
def get_ppt_file(date, chart,org_FAIl, data):
    """
    Create the power point file for a daily report in a bytes format.
    Args:
        date (datetime.date)    : User's selected date for the daily report.
        chart (plt)             : Figure (matplotlib format)
        org_FAIl (list)         : Org. Units that did not meet the daily goal.
        data (dataframe)        : Detailed table for org.unit, people reporting and
                                number of reports for each one.
    Returns:
        output (bytes)          : Power point report in bytes ready to be downloaded.
    """
    # Create empty ppt file
    prs = Presentation()
    top = Inches(1.3)
    left = Inches(0.5)
    height = Inches(5)
    # Layout for first and following slides
    first_slide_layout = prs.slide_layouts[0]
    second_layout=  prs.slide_layouts[1]
    # Adding slides
    slide_1 = prs.slides.add_slide(first_slide_layout)
    slide_2 = prs.slides.add_slide(second_layout)
    slide_3 = prs.slides.add_slide(second_layout)
    title_1 = slide_1.shapes.title
    subtitle_1 = slide_1.placeholders[1]
    # Formatting slides
    title_1.text = 'Hazards Report'
    subtitle_1.text = str(date)
    slide_2.shapes.title.text = 'Report'
    slide_3.shapes.title.text = 'Report Details - People'
    text_slide2 = slide_2.placeholders[1]
    text_slide2.text_frame.text = 'Did not meet the goal: {}'.format(org_FAIl)
    font = text_slide2.text_frame.paragraphs[0].runs[0].font
    font.size = Pt(10)

    #fig slide and more
    max_rows = 15
    counter = 0
    counter = 0
    for org_unit in np.unique(data['ORGANIZATIONAL UNIT']):
        data_unit = data[data['ORGANIZATIONAL UNIT']==org_unit]
        data_unit = data_unit[['Reported By','#R.']]
        data_unit = data_unit.sort_values('#R.', ascending=False).reset_index(drop=True)
        data_unit.loc[len(data_unit)] = ['------------------ TOTAL ------------------',data_unit['#R.'].sum()]
        no_tables = math.ceil(data_unit.shape[0]/max_rows)
        for i in range(no_tables):
            if counter == 3:
                slide_3 = prs.slides.add_slide(second_layout)
                slide_3.shapes.title.text = 'Report Details - People'
                counter = 0
            data_to_plot = data_unit[i*15:i*15+15]
            fig = return_df2_plot(data_to_plot, org_unit)
            image2 = BytesIO()
            df2img.save_dataframe(fig, image2)
            slide_3.shapes.add_picture(image2,Inches(3.2*(counter)+0.2), top, height= Inches(1.25)+Inches(data_to_plot.shape[0]/3))
            counter+=1
                    
    #Saving file
    image1 = BytesIO()
    chart.savefig(image1)
    slide_2.shapes.add_picture(image1,left, top, height=height)
    output = BytesIO()
    prs.save(output)
    return output   

def return_df2_plot(table, title):
    """
    Create the power point file for a daily report in a bytes format.
    Args:
        table (dataframe)    : Cropped table from detailed parent one.
        title (plt)          : Title for each cropped table, which corresponds to
                            the org. unit.
    Returns:
        fig (plt)           : Table and title converted to an image. Then, this
                            will be pasted in the PPT's slides.
    """
    fig = df2img.plot_dataframe(
        table,
        col_width=[0.05,0.29, 0.05],
        title=dict(
            font_color="darkred",
            font_family="Times New Roman",
            font_size=11,
            text=title,
        ),
        tbl_header=dict(
            align="left",
            fill_color="blue",
            font_color="white",
            font_size=8,
            line_color="darkslategray",
        ),
        tbl_cells=dict(
            align="left",
            line_color="darkslategray",
            font_color = 'black',
            font_size =7.5,
        ),
        row_fill_color=("#ffffff", "#d7d8d6"),
        fig_size=(195,80+20*table.shape[0])
    )
    return fig

def create_csv_bytes(hazards_file, text_file='dates.txt'):
    """
    Create a cvs file in bytes format to be downloaded for the daily report.
    Args:
        hazards_file (dataframe)    : Parent table that joined all of the day-to-day
                                    data.
    Returns:
        output.getvalue() (Bytes)  : Table to be download. This is assigned randomly from
                                    any date at hazards_file
        filename (str)             : Name of the file to be downloaded. The name will correspond
                                    to the next nate compared to the last date from hazards_file
    """
    text_file = open(text_file)
    split = text_file.readlines()[-1].split('_')
    text_file.close()
    year = split[3]
    month = split[4]
    day = str(split[5][:2])
    new_date = datetime.datetime.strptime(year+month+day,'%Y%b%d') + datetime.timedelta(days=1)
    filename = 'Events_SHEC_Hazards_'+'{}_{}_{}.csv'.format(new_date.year,calendar.month_abbr[new_date.month],new_date.day)
    # create the new df
    number = np.random.randint(1, len(np.unique(hazards_file['Event Date'])))
    date_selected = np.unique(hazards_file['Event Date'])[number]
    data_selected_date = hazards_file[hazards_file['Event Date']==date_selected]
    data_selected_date['Event Date'] = new_date
    output = BytesIO()
    data_selected_date.to_csv(output, index=False)
    return output.getvalue(), filename

def return_dict_months():
    """ Return a dictionary with the name of each month and its corresponding index"""
    dict_months = {'January': 1,
         'February': 2,
         'March': 3,
         'April': 4,
         'May': 5,
         'June': 6,
         'July': 7,
         'August': 8,
         'September': 9,
         'October': 10,
         'November': 11,
         'December': 12}
    return dict_months

def monthoptions(year_selected, year_now, month_now):
    """
    Create a list of options for the current year's month, i.e. if it is August 22,
    it will give you all the months between January and August.
    Args:
        year_selected (int)    : Year selected by the user
        year_now (int)         : Current year
        month_now (int)        : Current year's month
    Returns:
        ... (list)             : Corresponding months for last year's month.
    """
    total_months = return_dict_months()
    if year_selected< year_now:
        return total_months.keys()
    else:
        return [k for k,v in total_months.items() if v<month_now]
