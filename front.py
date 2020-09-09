import streamlit as st
import pptx
from pptx import Presentation
from pptx.util import Inches
#pylint: disable=E1120
import datetime
from calendar import monthrange
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib
import os
import back_
import plotly.express as px
import plotly.graph_objects as go
import colorama
from colorama import Fore, Style
from plotly.subplots import make_subplots
from IPython.display import Image
from PIL import Image
from io import BytesIO
from io import StringIO 
import base64


#ORG SELECTED FOR REPORTS:
selected_day = {'COMMUNITY RELATIONS': 10, 'CONCENTRATOR PLANT':118, 'ENVIRONMENT':6, 'LAS BAMBAS - PROJECT DELIVERY': 18, 
            'LOGISTICS': 23, 'MAINTENANCE': 100, 'MINE OPERATIONS': 270, 'PLANNING, INVENTORY & WAREHOUSE': 65, 'SAFETY & HEALTH':7,
             'TECHNICAL SERVICES': 37}
#TYPE OF INCIDENTS
recordables = ['FI', 'LTI', 'MTI', 'RWI']
national = ['FI', 'LTI']
all_injuries = ['FI', 'LTI', 'MTI', 'RWI', 'FAI']
lost_time = ['LTI']
fatal = ['FI']
rest = ['RWI']
mti = ['MTI']
fai = ['FAI']


hazards_dtb= pd.read_excel("hazard_total.xlsx").reset_index(drop = True)
hazards_dtb.iloc[:,1] = pd.to_datetime(hazards_dtb.iloc[:,1]).dt.date
Main = back_.Carga_Reports(hazards_dtb, "new_accidentes.xlsx", "horas_hombre.xlsx")
def texto(valor):
    if valor < 1:
        return "x"
    else:
        return ""




def main():
    st.title("Safety KPI Dashboard")
    menu = ['Daily Report', 'Weekly Report', 'Monthly Report']
    choice = st.sidebar.selectbox('Report Type', menu)
    today = datetime.date.today()
    if choice == 'Daily Report':    
        day_daily = st.sidebar.date_input("Select a day:", datetime.date.today())
        if day_daily > today:
            st.sidebar.error("Change dates")
        else:
            go_daily = st.sidebar.checkbox("View Report:")
            if go_daily:
                try:
                    st.success("Report date: {}".format(day_daily))                
                    st.plotly_chart(get_daily_graph(day_daily = day_daily))
            
                    st.success("People reporting on {}". format(day_daily))
                    st.write(get_daily_dataframe())   
                    #How to download this data???????
                    st.markdown(get_link_download(), unsafe_allow_html= True) 
                    
                except:
                    st.write("We do not have data for this day")

    if choice == 'Weekly Report':
        initial_date = st.sidebar.date_input("From:", datetime.date.today())
        end_date = st.sidebar.date_input("To:", datetime.date.today())
        if initial_date > today or end_date > today or initial_date > end_date:
            st.sidebar.error("Change dates")
        else:
            go_weekly = st.sidebar.checkbox("View Report:")
            if go_weekly:
                st.header(("Report from {} to {}").format(initial_date, end_date))
                weekly = back_.Weekly_data(hazards_dtb, "new_accidentes.xlsx", "horas_hombre.xlsx", initial_date, end_date)
                info_week = get_weekly_part1(initial_date, end_date, Main, all_injuries, recordables, lost_time, weekly)
                #First dataframe
                st.subheader("OSHA Monitoring Indicators")
                st.dataframe(info_week[0].style.format("{:,.2f}"))
                st.dataframe(info_week[1].style.format("{:,.2f}"))
                #Third part
                st.subheader("National Monitoring Indicators")
                st.dataframe(info_week[2].style.format("{:.2f}"))
                #Fourth Part
                st.subheader("Injuries and Illness")
                st.dataframe(info_week[3].style.format("{:.0f}"))
                #fifth part
                st.subheader("Incidentes of the week")
                st.dataframe(info_week[4])
                #sixth part
                st.subheader("MLB Safety Performance")
                st.success("Year = YTD KPI and Month = MTD KPI")
                st.plotly_chart(info_week[5])
                #Seventh part
                #we can add the intersection thing!!!!!!
                st.plotly_chart(info_week[6])

    if choice == 'Monthly Report':
        year = st.sidebar.selectbox("Year: ", range(2018, 2024))
        dict_months = {'January': 1,
         'February': 2,
         'March': 3,
         'April': 4,
         'May': 5,
         'June': 6,
         'July': 7,
         'August': 8,
         'September': 9,
         'October': 10,
         'November': 11,
         'December': 12}
        month = st.sidebar.selectbox("Month: ", list(dict_months.keys()))
        start =  datetime.datetime.strptime(str(year)+"-"+str(dict_months[month])+"-01 00:00:00", "%Y-%m-%d %H:%M:%S")
        end = datetime.datetime.strptime(str(year)+"-"+str(dict_months[month])+ "-" + str(monthrange(year, dict_months[month])[1])+" 00:00:00", "%Y-%m-%d %H:%M:%S")
        
        if year > datetime.datetime.now().year or (year >= datetime.datetime.now().year and dict_months[month] > datetime.datetime.now().month):
            st.sidebar.error("Change dates")
        else:
            go_monthly = st.sidebar.checkbox("View Report: ")
            if go_monthly:
                st.header(("Report of {}").format(start.strftime("%B-%Y")))
                st.subheader("OSHA Monitoring Indicators")
                weekly = back_.Weekly_data(hazards_dtb, "new_accidentes.xlsx", "horas_hombre.xlsx", start, end)
                lista = [end.strftime("%B-%Y")] + [_ for _ in range(end.year, min(Main.file_incidents['FECHA'].dt.year)-1 , -1)]
                #First Part
                first_part = pd.DataFrame(columns = lista)
                first_part.loc['All Injuries Frequency Rate (AIFR)'] = [weekly.kpi_mtd(selected = all_injuries)] + weekly.kpi_ytd(selected= all_injuries)
                first_part.loc['Total Recordable Injuries Freq. Rate (TRIFR)'] = [weekly.kpi_mtd(selected = recordables)] + weekly.kpi_ytd(selected= recordables)
                
                #ISSSUEEE!!!!!!!!1
                first_part.loc['Lost Time Injuries Freq. Rate (LTIFR)'] = [weekly.kpi_mtd(selected = lost_time)] + weekly.kpi_ytd(selected= lost_time)
                st.dataframe(first_part.style.format("{:.2f}"))

                
                #second part of the first part
                second_part = pd.DataFrame(columns = [""], index = ["TRIFR 12MMA (to {})".format(weekly.to_dataframe(moving=12).strftime("%Y-%B-%d")), \
                    "Man-Hoursworked no LTI (Last one was on: {})".format(weekly.last_lti_(selected = 'LTI').strftime("%Y-%B-%d")), "Lost Days (YTD)"])
                second_part.iloc[:,0] = [weekly.int_kpi_mma(selected = recordables, moving = 12), weekly.manhoursworkerd_nolti(selected = 'LTI'), weekly.number_lost_day()]
                st.dataframe(second_part.style.format("{:,.2f}"))
                
                #Third part
                st.subheader("National Monitoring Indicators")
                third_part = pd.DataFrame(columns = lista)
                #index = [, 'Severity Index (SI)','Accidentability Index (AI)']
                third_part.loc['Frequency Index (FI)'] = [weekly.kpi_mtd(selected = national, horas = "todas")] + weekly.kpi_ytd(selected = national, horas = "todas")
                third_part.loc['Severity Index (SI)'] = weekly.national_si_kpi()
                third_part = third_part.fillna(0)

                third_part.loc['Accidentability Index (AI)'] = third_part.loc['Frequency Index (FI)'] * third_part.loc['Severity Index (SI)'] / 1000

                st.dataframe(third_part.style.format("{:.2f}"))

            
                #Fourth Part
                st.subheader("Injuries and Illness")
                fourth_part = pd.DataFrame(columns = lista)
                fourth_part.loc['A. Fatalities'] = [weekly.int_kpi_mtd(selected = fatal)] + weekly.int_kpi_ytd(selected = fatal)
                fourth_part.loc['B. Lost Time Inj.'] =  [weekly.int_kpi_mtd(selected = lost_time)] + weekly.int_kpi_ytd(selected = lost_time)
                fourth_part.loc['C. Rest. Work Inj.'] =  [weekly.int_kpi_mtd(selected = rest)] + weekly.int_kpi_ytd(selected = rest)
                fourth_part.loc['D. Medic. Tre. Inj.'] =   [weekly.int_kpi_mtd(selected = mti)] + weekly.int_kpi_ytd(selected = mti)

                fourth_part.loc['E. Recordable Injuries'] = fourth_part.sum(axis = 0) 
                
                fourth_part.loc["F. First Aid Inj."] = [weekly.int_kpi_mtd(selected = fai)] + weekly.int_kpi_ytd(selected = fai)
                fourth_part.loc["G. Total Injuries:"] = fourth_part.loc[['E. Recordable Injuries', "F. First Aid Inj."]].sum(axis = 0) 
                st.dataframe(fourth_part.style.format("{:.0f}"))

                #fifth part
                st.subheader("Incidentes of the month")
                st.dataframe(weekly.incidents())

                #sixth part
                st.subheader("Monthly - MLB Safety Performance")
                
                cols_c = fourth_part.columns.to_list()[::-1]
                index_c = ["A. Fatalities", "B. Lost Time Inj.","C. Rest. Work Inj.","D. Medic. Tre. Inj.","F. First Aid Inj."]
                up = pd.DataFrame(fourth_part, columns = cols_c, index = index_c)
                low = pd.DataFrame(first_part, index = ['Total Recordable Injuries Freq. Rate (TRIFR)', 'Lost Time Injuries Freq. Rate (LTIFR)', 'All Injuries Frequency Rate (AIFR)'], columns = cols_c)
                last = pd.concat([up, low])
                last.index = ['FI', 'LTI', 'RWI', 'MTI', 'FAI', 'TRIFR', 'LTIFR', 'AIFR']
                last = last.T
                #chart
                fig = go.Figure()
                color1 = ['#800000', '#3366ff', '#ff9933', '#00ff00', '#ffff00']
                color2 = ['#003300','#ff00ff' ]
                for column, color in zip(last.columns[:5], color1):
                    fig.add_trace(go.Bar(x = [str(x) for x in last.index], y = last[column], marker_color = color, name = column, yaxis = 'y1'))
                    
                for column, color in zip(last.columns[[5,7]], color2):
                    x_axis = [str(x) for x in last.index]
                    y_axis = [round(elem, 2) for elem in last[column]]
                    fig.add_trace(go.Scatter(x = x_axis, y = y_axis, marker_color = color, name = column, yaxis = 'y2', \
                    mode = 'lines+markers+text'))#, text = text, textfont = dict(color = color)))
                    for x,y in zip(x_axis, y_axis):
                        fig.add_annotation(x = x, y = y, yref = 'y2',text = y,  font = dict(color = "white"), align = "center", bgcolor = color)
                    
                fig.update_layout(barmode = 'stack', xaxis = {'type' : 'category'}, yaxis = go.layout.YAxis(), yaxis2 = go.layout.YAxis(side = 'right', overlaying = 'y1'), legend = dict(orientation = 'h', yanchor = 'bottom', y = 1.02, xanchor = 'right', x = 1),\
                    title = "Injuries & KPIs / Date", xaxis_title = 'Date', yaxis_title = 'Injuries Frequency', yaxis2_title = "TRIFR / AIFR")
                st.plotly_chart(fig)

                #Seventh part
                x_axis = weekly.names_mma_charts(selected = recordables)
                fig = go.Figure()
                fig.add_trace(go.Scatter(x = x_axis, y = weekly.mma_charts(selected = recordables, moving =12), mode = 'lines+markers', name = 'TRIFR-12MMA', marker_color = '#000000'))
                fig.add_trace(go.Scatter(x = x_axis, y = weekly.mma_charts(selected = recordables, moving =3), mode = 'lines+markers', name = 'TRIFR-3MMA', marker_color = '#c4bcbc'))
                fig.update_layout(legend = dict(yanchor = 'top', y = 0.99, xanchor = 'left', x = 0.01), title = "TRIFR 12MMA & 3MMA / Month", yaxis_title = "TRIFR", xaxis_title = "Date")
                #we can add the intersection thing!!!!!!
                st.plotly_chart(fig)


# Daily Report - chart (this is Hazrds KPI & Hazrd. Freq. vs Org. Unit)     
def get_daily_graph(day_daily):
    global daily
    daily = back_.Daily_report(hazards_dtb, "new_accidentes.xlsx", "horas_hombre.xlsx", day_daily, selected_daily= selected_day)
    daily_ = back_.Daily_report(hazards_dtb, "new_accidentes.xlsx", "horas_hombre.xlsx", day_daily - datetime.timedelta(days = 1), selected_daily= selected_day)
    fig = go.Figure(go.Bar(x = daily_.graphs_daily().iloc[:,0], y = daily_.graphs_daily().iloc[:,1], yaxis = 'y1', marker_color = "#AAF0D1",\
        name = str(day_daily- datetime.timedelta(days = 1))))
    fig.add_trace(go.Bar(x = daily.graphs_daily().iloc[:,0], y = daily.graphs_daily().iloc[:,1], yaxis = 'y1', marker_color = "#50BFE6",\
        name = str(day_daily)))
    fig.add_trace(go.Scatter(x = daily.graphs_daily().iloc[:,0], y = daily.graphs_daily().iloc[:,3], mode = 'markers+text', \
        showlegend = False, marker = dict(color = "#50BFE6", size = 20, line = dict(color = '#FFFFFF', width = 3)), text = list(map(texto,daily.graphs_daily().iloc[:,3])), textfont = dict(color = "#800000", size = 32),\
        yaxis = 'y2', name = ""))
    fig.add_trace(go.Scatter(x = daily_.graphs_daily().iloc[:,0], y = daily_.graphs_daily().iloc[:,3], 
    mode = 'lines+markers', showlegend= False, marker_color = "#AAF0D1", yaxis = 'y2', name = ""))
    fig.update_layout(barmode = 'stack', yaxis = go.layout.YAxis(), yaxis2 = go.layout.YAxis(side = 'right', overlaying = 'y'), legend = dict(orientation = 'h', yanchor = 'bottom', xanchor = 'right', y = 1.1, x = 1),\
        xaxis_title = '<b>Organizational Units</b>', yaxis_title = '<b>Hazards Frequency</b>', yaxis2_title = '<b>Hazard KPI</b>')
    return fig

# Daily Rport - A detailed table of whom has reported on the day basis
def get_daily_dataframe():
    return daily.data_daily()  
# Daily Report - A link to download the daily report in pptx format but this is still on dev.
def get_link_download():                      
    prs = Presentation()
    first_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(first_slide_layout)
    slide.shapes.title.text = "Report date"
    slide.placeholders[1].text = 'This is 2nd way'
    
    second_slide_layout = prs.slide_layouts[1]
    slide1 = prs.slides.add_slide(second_slide_layout)
    #slide1.insert_picture(fig.to_image(format = 'JPEG'))
    #img_bytes = fig.to_image(format="png")
    image1 = BytesIO()
    fig.write_image(image1)
    image1_toput = Image.frombytes(image1)
    slide1.shapes.title.text = "Report date: {}".format(day_daily)
    slide1.shapes.add_picute(image1_toput, Inches(1), Inches(1))
    #slide1.placeholders[1].(Image(img_bytes), Inches(1), Inches(1))
    #img_bytes = fig.to_image(format="png", width = 100, height = 100, scale = 2)
    
    #out_file = base64.b64encode(out_file)
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    nuevo = base64.b64encode(output.getbuffer()).decode()
    #b = a.write(output.getvalue())


    #title.text = "Report date: {}".format(day_daily)
    #placeholder = slide.placeholders[0]
    #placeholder.insert_picture(img.to_image(format = 'JPEG'))
    #buffered = img.to_image(format = 'JPEG')
    href = f'<a href="data:file/pptx;base64, {nuevo}">Download report in ptt format</a>'
    return href   


def get_weekly_part1(initial_date, end_date, Main, all_injuries, recordables, lost_time, weekly):
    lista = [end_date.strftime("%B-%Y")] + [_ for _ in range(end_date.year, min(Main.file_incidents['FECHA'].dt.year)-1 , -1)]
    #OSHA MONITORING - first table
    first_part = pd.DataFrame(columns = lista)
    first_part.loc['All Injuries Frequency Rate (AIFR)'] = [weekly.kpi_mtd(selected = all_injuries)] + weekly.kpi_ytd(selected= all_injuries)
    first_part.loc['Total Recordable Injuries Freq. Rate (TRIFR)'] = [weekly.kpi_mtd(selected = recordables)] + weekly.kpi_ytd(selected= recordables)
    first_part.loc['Lost Time Injuries Freq. Rate (LTIFR)'] = [weekly.kpi_mtd(selected = lost_time)] + weekly.kpi_ytd(selected= lost_time)

    #Belongs to the first table
    second_part = pd.DataFrame(columns = [""], index = ["TRIFR 12MMA (to {})".format(weekly.to_dataframe(moving=12).strftime("%Y-%B-%d")), \
        "Man-Hoursworked no LTI (Last one was on: {})".format(weekly.last_lti_(selected = 'LTI').strftime("%Y-%B-%d")), "Lost Days (YTD)"])
    second_part.iloc[:,0] = [weekly.int_kpi_mma(selected = recordables, moving = 12), weekly.manhoursworkerd_nolti(selected = 'LTI'), weekly.number_lost_day()]
    
    # National monitoring indicators - third table
    third_part = pd.DataFrame(columns = lista)
        #index = [, 'Severity Index (SI)','Accidentability Index (AI)']
    third_part.loc['Frequency Index (FI)'] = [weekly.kpi_mtd(selected = national, horas = "todas")] + weekly.kpi_ytd(selected = national, horas = "todas")
    third_part.loc['Severity Index (SI)'] = weekly.national_si_kpi()
    third_part = third_part.fillna(0)
    third_part.loc['Accidentability Index (AI)'] = third_part.loc['Frequency Index (FI)'] * third_part.loc['Severity Index (SI)'] / 1000

    # Injuries and Illnesses - fourth table
    fourth_part = pd.DataFrame(columns = [initial_date.strftime("%m-%d") + " to " + end_date.strftime("%m-%d")]+ lista)
    fourth_part.loc['A. Fatalities'] = [weekly.int_kpi_wtd(selected = fatal)] + [weekly.int_kpi_mtd(selected = fatal)] + weekly.int_kpi_ytd(selected = fatal)
    fourth_part.loc['B. Lost Time Inj.'] = [weekly.int_kpi_wtd(selected = lost_time)] + [weekly.int_kpi_mtd(selected = lost_time)] + weekly.int_kpi_ytd(selected = lost_time)
    fourth_part.loc['C. Rest. Work Inj.'] = [weekly.int_kpi_wtd(selected = rest)] + [weekly.int_kpi_mtd(selected = rest)] + weekly.int_kpi_ytd(selected = rest)
    fourth_part.loc['D. Medic. Tre. Inj.'] = [weekly.int_kpi_wtd(selected = mti)] + [weekly.int_kpi_mtd(selected = mti)] + weekly.int_kpi_ytd(selected = mti)

    fourth_part.loc['E. Recordable Injuries'] = fourth_part.sum(axis = 0) 
    
    fourth_part.loc["F. First Aid Inj."] = [weekly.int_kpi_wtd(selected = fai)] + [weekly.int_kpi_mtd(selected = fai)] + weekly.int_kpi_ytd(selected = fai)
    fourth_part.loc["G. Total Injuries:"] = fourth_part.loc[['E. Recordable Injuries', "F. First Aid Inj."]].sum(axis = 0) 
    
    # Incidents table: Fifth part
    if len(weekly.incidents()) <= 0:
        fifth_part = "Nice, we did not have incidents!!"
    else:
        fifth_part = weekly.incidents()

    #TRIFR, AIFR CHART: Sixth part
    cols_c = fourth_part.columns.to_list()[::-1][:-1]
    index_c = ["A. Fatalities", "B. Lost Time Inj.","C. Rest. Work Inj.","D. Medic. Tre. Inj.","F. First Aid Inj."]
    up = pd.DataFrame(fourth_part, columns = cols_c, index = index_c)
    low = pd.DataFrame(first_part, index = ['Total Recordable Injuries Freq. Rate (TRIFR)', 'Lost Time Injuries Freq. Rate (LTIFR)', 'All Injuries Frequency Rate (AIFR)'], columns = cols_c)
    last = pd.concat([up, low])
    last.index = ['FI', 'LTI', 'RWI', 'MTI', 'FAI', 'TRIFR', 'LTIFR', 'AIFR']
    last = last.T
    fig = go.Figure()
    color1 = ['#800000', '#3366ff', '#ff9933', '#00ff00', '#ffff00']
    color2 = ['#003300','#ff00ff' ]
    for column, color in zip(last.columns[:5], color1):
        fig.add_trace(go.Bar(x = [str(x) for x in last.index], y = last[column], marker_color = color, name = column, yaxis = 'y1'))
        
    for column, color in zip(last.columns[[5,7]], color2):
        x_axis = [str(x) for x in last.index]
        y_axis = [round(elem, 2) for elem in last[column]]
        fig.add_trace(go.Scatter(x = x_axis, y = y_axis, marker_color = color, name = column, yaxis = 'y2', \
        mode = 'lines+markers+text'))#, text = text, textfont = dict(color = color)))
        for x,y in zip(x_axis, y_axis):
            fig.add_annotation(x = x, y = y, yref = 'y2',text = y,  font = dict(color = "white"), align = "center", bgcolor = color)
        
    fig.update_layout(barmode = 'stack', xaxis = {'type' : 'category'}, yaxis = go.layout.YAxis(), yaxis2 = go.layout.YAxis(side = 'right', overlaying = 'y1'), legend = dict(orientation = 'h'))
    
    x_axis = weekly.names_mma_charts(selected = recordables)
    fig2 = go.Figure()
    fig2.add_trace(go.Scatter(x = x_axis, y = weekly.mma_charts(selected = recordables, moving =12), mode = 'lines+markers', name = 'TRIFR-12MMA', marker_color = '#000000'))
    fig2.add_trace(go.Scatter(x = x_axis, y = weekly.mma_charts(selected = recordables, moving =3), mode = 'lines+markers', name = 'TRIFR-3MMA', marker_color = '#c4bcbc'))
    fig2.update_layout(legend = dict(orientation = 'h'))

    return [first_part, second_part, third_part, fourth_part, fifth_part, fig, fig2]

if __name__ == '__main__':
    main()